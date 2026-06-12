#!/usr/bin/env python3
"""
pptx_to_html — Convert a PowerPoint dashboard slide into clean, responsive HTML.

Approach:
  1. Parse shapes with python-pptx (geometry, fills, borders, text runs, images,
     tables, connectors). Resolve theme colors (with the master color map) and
     inherited placeholder styles from the slide layout / master.
  2. Build a containment tree from shape geometry (which shapes sit inside which).
  3. Group siblings into rows (vertical bands) and columns (horizontal order),
     emitting flexbox/grid layout instead of absolute positioning.
  4. Emit semantic HTML + a CSS file with variables for every color and font,
     preserving exact colors, font sizes, weights, and spacing.

Usage:
  python -m ai_doc_converter input.pptx [-o output_dir] [--slide N]

Output:
  output_dir/index.html, styles.css, assets/  (extracted images)
"""

import argparse
import os
import re
import sys
from html import escape

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu

EMU_PER_PX = 9525          # 96 px per inch
ROW_OVERLAP = 0.5          # min vertical overlap ratio to share a row
CONTAIN_TOL_PX = 2         # tolerance when testing geometric containment
MIN_BAND_PX = 8            # effective height for zero-height shapes (lines)
LABEL_SNAP_PX = 40         # max distance to attach a text label to a connector

CONNECTOR_PRST = re.compile(
    r'prst="(line|straightConnector\d*|bentConnector\d*|curvedConnector\d*)"')

ALIGN_MAP = {PP_ALIGN.LEFT: "left", PP_ALIGN.CENTER: "center",
             PP_ALIGN.RIGHT: "right", PP_ALIGN.JUSTIFY: "justify"}
ALGN_ATTR = {"l": "left", "ctr": "center", "r": "right", "just": "justify"}
ANCHOR_MAP = {MSO_ANCHOR.TOP: "flex-start", MSO_ANCHOR.MIDDLE: "center",
              MSO_ANCHOR.BOTTOM: "flex-end"}

# MSO_THEME_COLOR enum name -> theme/colorMap key
THEME_ENUM_KEY = {
    "DARK_1": "dk1", "DARK_2": "dk2", "LIGHT_1": "lt1", "LIGHT_2": "lt2",
    "ACCENT_1": "accent1", "ACCENT_2": "accent2", "ACCENT_3": "accent3",
    "ACCENT_4": "accent4", "ACCENT_5": "accent5", "ACCENT_6": "accent6",
    "HYPERLINK": "hlink", "FOLLOWED_HYPERLINK": "folHlink",
    "TEXT_1": "tx1", "TEXT_2": "tx2", "BACKGROUND_1": "bg1", "BACKGROUND_2": "bg2",
}

# Map non-web-safe PowerPoint fonts to sensible fallback stacks.
FONT_FALLBACKS = {
    "Gadugi": "'Gadugi', 'Segoe UI', 'Helvetica Neue', Arial, sans-serif",
    "Calibri": "'Calibri', 'Segoe UI', Arial, sans-serif",
    "Calibri Light": "'Calibri Light', 'Segoe UI Light', 'Segoe UI', sans-serif",
    "Georgia": "Georgia, 'Times New Roman', serif",
}

SKIP_PLACEHOLDERS = ("SLIDE_NUMBER", "FOOTER", "DATE")


def px(emu):
    return round(emu / EMU_PER_PX, 1)


# ------------------------------------------------------------------ theming

def scheme_hex(name, theme):
    """Resolve a schemeClr name (tx1, bg1, dk1, accent1...) through the
    master color map to a theme hex color."""
    name = theme.get("clrmap", {}).get(name, name)
    return theme.get("colors", {}).get(name)


def color_of(color_format, theme):
    """Resolve a python-pptx ColorFormat to a hex string, or None."""
    try:
        if color_format.type is None:
            return None
        if str(color_format.type).startswith("MSO_THEME_COLOR"):
            key = THEME_ENUM_KEY.get(str(color_format.theme_color).split()[0])
            return scheme_hex(key, theme) if key else None
        return f"#{color_format.rgb}"
    except (AttributeError, TypeError):
        return None


def _parse_color(segment, theme):
    m = re.search(r'<a:srgbClr val="([0-9A-Fa-f]{6})"', segment)
    if m:
        return "#" + m.group(1).upper()
    m = re.search(r'<a:schemeClr val="(\w+)"', segment)
    if m:
        return scheme_hex(m.group(1), theme)
    return None


def read_theme(prs):
    """Extract theme palette, fonts and the master color map."""
    theme = {"colors": {}, "clrmap": {}, "major_font": None, "minor_font": None}
    master = prs.slide_masters[0]
    m = re.search(r'<p:clrMap([^/>]*)/?>', master.element.xml)
    if m:
        theme["clrmap"] = dict(re.findall(r'(\w+)="(\w+)"', m.group(1)))
    for rel in master.part.rels.values():
        if "theme" in rel.reltype:
            xml = rel.target_part.blob.decode("utf-8", errors="ignore")
            for name in ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2",
                         "accent3", "accent4", "accent5", "accent6",
                         "hlink", "folHlink"):
                mm = re.search(rf'<a:{name}>.*?(?:val|lastClr)="([0-9A-Fa-f]{{6}})"', xml, re.S)
                if mm:
                    theme["colors"][name] = "#" + mm.group(1).upper()
            mm = re.search(r'<a:majorFont><a:latin typeface="([^"]*)"', xml)
            theme["major_font"] = mm.group(1) if mm else "Calibri Light"
            mm = re.search(r'<a:minorFont><a:latin typeface="([^"]*)"', xml)
            theme["minor_font"] = mm.group(1) if mm else "Calibri"
    return theme


def slide_background(slide, theme):
    """Resolve the slide background color from slide -> layout -> master."""
    for part in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        m = re.search(r'<p:bg>.*?</p:bg>', part.element.xml, re.S)
        if m:
            c = _parse_color(m.group(0), theme)
            if c:
                return c
    return "#FFFFFF"


# ------------------------------------------------- placeholder inheritance

def parse_lvl_style(lststyle_xml, level, theme):
    """Parse one list-style level (lvl1pPr...) into a style dict."""
    m = re.search(rf'<a:lvl{level + 1}pPr([^>]*)>(.*?)</a:lvl{level + 1}pPr>',
                  lststyle_xml, re.S)
    p_attrs, seg = (m.group(1), m.group(2)) if m else ("", lststyle_xml)
    out = {}
    am = re.search(r'algn="(\w+)"', p_attrs)
    if am and ALGN_ATTR.get(am.group(1)):
        out["align"] = ALGN_ATTR[am.group(1)]
    dm = re.search(r'<a:defRPr([^>]*)(?:/>|>(.*?)</a:defRPr>)', seg, re.S)
    if not dm:
        return out
    attrs, body = dm.group(1), dm.group(2) or ""
    sm = re.search(r'\bsz="(\d+)"', attrs)
    if sm:
        out["size"] = round(int(sm.group(1)) / 100, 1)
    if re.search(r'\bb="1"', attrs):
        out["bold"] = True
    if re.search(r'\bi="1"', attrs):
        out["italic"] = True
    c = _parse_color(body, theme)
    if c:
        out["color"] = c
    fm = re.search(r'<a:latin typeface="([^"]+)"', body)
    if fm:
        out["font"] = fm.group(1)
    return out


def placeholder_fallback(shape, layout, master, theme, level=0):
    """Inherited run style for a placeholder: layout ph -> master ph ->
    master txStyles."""
    if not getattr(shape, "is_placeholder", False):
        return {}
    try:
        pf = shape.placeholder_format
        idx, typ = pf.idx, pf.type
    except (AttributeError, ValueError):
        return {}
    def norm(t):
        return "title" if "TITLE" in str(t or "").upper() else "body"

    def find_ph(src, by="idx"):
        by_idx = by_type = None
        for ph in src.placeholders:
            try:
                pf2 = ph.placeholder_format
                if pf2.idx == idx and by_idx is None:
                    by_idx = ph
                if norm(pf2.type) == norm(typ) and by_type is None:
                    by_type = ph
            except (AttributeError, ValueError):
                continue
        return (by_idx or by_type) if by == "idx" else (by_type or by_idx)

    sty = {}
    # layout: same placeholder (idx match); master: prototype by type
    for src, mode in ((layout, "idx"), (master, "type")):
        if src is None:
            continue
        ph = find_ph(src, mode)
        if ph is None:
            continue
        m = re.search(r'<a:lstStyle>(.*?)</a:lstStyle>', ph._element.xml, re.S)
        if not m:
            continue
        for k, v in parse_lvl_style(m.group(1), level, theme).items():
            sty.setdefault(k, v)
        if "size" in sty and "color" in sty:
            return sty
    if master is not None:
        tname = "titleStyle" if "TITLE" in str(typ or "").upper() else "bodyStyle"
        m = re.search(rf'<p:{tname}>(.*?)</p:{tname}>', master.element.xml, re.S)
        if m:
            for k, v in parse_lvl_style(m.group(1), level, theme).items():
                sty.setdefault(k, v)
    return sty


# ---------------------------------------------------------------- extraction

def detect_connector(shape, d):
    """Detect line/connector shapes: flow arrows or plain divider lines."""
    try:
        xml = shape._element.spPr.xml
    except AttributeError:
        return False
    if not CONNECTOR_PRST.search(xml):
        return False
    d["dir"] = "h" if d["w"] >= d["h"] else "v"
    m = re.search(r'<a:ln[^>]*>.*?<a:srgbClr val="([0-9A-Fa-f]{6})"', xml, re.S)
    if m:
        d["stroke"] = "#" + m.group(1).upper()
    else:
        m = re.search(r'<a:ln[^>]*>.*?<a:schemeClr val="(\w+)"', xml, re.S)
        d["stroke"] = (scheme_hex(m.group(1), d["_theme"]) if m else None) or "#444444"
    wm = re.search(r'<a:ln w="(\d+)"', xml)
    d["stroke_w"] = max(1, round(int(wm.group(1)) / EMU_PER_PX)) if wm else 1
    d["arrow_start"] = bool(re.search(r'<a:headEnd type="(?!none)', xml))
    d["arrow_end"] = bool(re.search(r'<a:tailEnd type="(?!none)', xml))
    d["kind"] = "connector" if (d["arrow_start"] or d["arrow_end"]) else "divider"
    return True


def run_style(run, theme):
    f = run.font
    return {
        "font": f.name,
        "size": round(f.size.pt, 1) if f.size else None,
        "bold": f.bold,
        "italic": f.italic,
        "color": color_of(f.color, theme),
    }


def extract_text(shape, theme, fallback=None):
    fb = fallback or {}
    paras = []
    for p in shape.text_frame.paragraphs:
        runs = []
        for r in p.runs:
            if not r.text:
                continue
            rs = run_style(r, theme)
            for key in ("font", "size", "bold", "italic", "color"):
                if rs[key] is None:
                    rs[key] = fb.get(key)
            runs.append({"text": r.text, **rs})
        ls = p.line_spacing
        if ls is not None and not isinstance(ls, float):
            ls = round(ls.pt / 12.0, 2)  # approximate multiple
        paras.append({
            "runs": runs,
            "align": ALIGN_MAP.get(p.alignment) or fb.get("align"),
            "line_spacing": ls,
            "space_after": round(p.space_after.pt, 1) if p.space_after else None,
        })
    return paras


def extract_table(shape, theme):
    rows = []
    for r in shape.table.rows:
        row = []
        for c in r.cells:
            row.append({
                "paras": [{"runs": [{"text": run.text, **run_style(run, theme)}
                                    for run in p.runs],
                           "align": ALIGN_MAP.get(p.alignment),
                           "line_spacing": None, "space_after": None}
                          for p in c.text_frame.paragraphs],
                "fill": color_of(c.fill.fore_color, theme)
                        if c.fill.type is not None and c.fill.type == 1 else None,
            })
        rows.append(row)
    return rows


def extract_shape(shape, theme, assets_dir, asset_list, layout=None, master=None):
    if getattr(shape, "is_placeholder", False):
        try:
            if any(k in str(shape.placeholder_format.type or "")
                   for k in SKIP_PLACEHOLDERS):
                return None   # page numbers / footers / dates: skip on the web
        except (AttributeError, ValueError):
            pass

    d = {
        "id": shape.shape_id, "name": shape.name,
        "x": px(shape.left), "y": px(shape.top),
        "w": px(shape.width), "h": px(shape.height),
        "kind": "box", "fill": None, "border": None, "radius": 0,
        "paras": None, "table": None, "img": None, "anchor": None,
        "_theme": theme,
    }
    d["x2"], d["y2"] = d["x"] + d["w"], d["y"] + d["h"]

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        d["kind"] = "img"
        image = shape.image
        fname = f"image{len(asset_list) + 1}.{image.ext}"
        with open(os.path.join(assets_dir, fname), "wb") as fh:
            fh.write(image.blob)
        asset_list.append(fname)
        d["img"] = f"assets/{fname}"
        return d

    if shape.has_table:
        d["kind"] = "table"
        d["table"] = extract_table(shape, theme)
        return d

    if detect_connector(shape, d):
        return d

    try:
        if shape.fill.type is not None and shape.fill.type == 1:  # solid
            d["fill"] = color_of(shape.fill.fore_color, theme)
    except (AttributeError, TypeError):
        pass
    try:
        line = shape.line
        if line.fill.type is not None and line.fill.type == 1 and line.width:
            d["border"] = {"color": color_of(line.color, theme),
                           "width": px(line.width)}
    except (AttributeError, TypeError):
        pass
    # rounded corners (default PowerPoint adjustment is 16.67% when unset)
    try:
        if "round" in str(shape.auto_shape_type or "").lower():
            adjs = list(shape.adjustments)
            frac = adjs[0] if adjs else 0.1667
            d["radius"] = round(min(frac, 0.5) * min(d["w"], d["h"]), 1)
    except (AttributeError, ValueError, IndexError):
        pass

    if shape.has_text_frame and shape.text_frame.text.strip():
        d["kind"] = "text"
        fb = placeholder_fallback(shape, layout, master, theme)
        d["paras"] = extract_text(shape, theme, fb)
        d["anchor"] = ANCHOR_MAP.get(shape.text_frame.vertical_anchor)
    return d


def flatten(shapes, theme, assets_dir, asset_list, layout=None, master=None):
    out = []
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            out.extend(flatten(sh.shapes, theme, assets_dir, asset_list, layout, master))
        else:
            try:
                d = extract_shape(sh, theme, assets_dir, asset_list, layout, master)
                if d:
                    out.append(d)
            except Exception as e:
                print(f"  warn: skipped shape {sh.shape_id} ({e})", file=sys.stderr)
    return out


# ------------------------------------------------------------------- layout

def contains(a, b):
    t = CONTAIN_TOL_PX
    return (a is not b
            and a["x"] - t <= b["x"] and a["y"] - t <= b["y"]
            and a["x2"] + t >= b["x2"] and a["y2"] + t >= b["y2"]
            and a["w"] * a["h"] > b["w"] * b["h"])


def build_tree(shapes):
    """Assign each shape to its smallest geometric container."""
    for s in shapes:
        s["children"] = []
    roots = []
    for s in shapes:
        candidates = [c for c in shapes if contains(c, s)]
        if candidates:
            parent = min(candidates, key=lambda c: c["w"] * c["h"])
            parent["children"].append(s)
        else:
            roots.append(s)
    return roots


def group_rows(children):
    """Group sibling shapes into horizontal rows by vertical overlap."""
    rows = []
    for s in sorted(children, key=lambda c: (c["y"], c["x"])):
        placed = False
        h_eff = max(s["h"], MIN_BAND_PX)   # zero-height lines still need a band
        for row in rows:
            ry1 = min(c["y"] for c in row)
            ry2 = max(c["y2"] for c in row)
            ov = min(ry2, s["y"] + h_eff) - max(ry1, s["y"])
            if ov > ROW_OVERLAP * min(h_eff, max(ry2 - ry1, MIN_BAND_PX)):
                row.append(s)
                placed = True
                break
        if not placed:
            rows.append([s])
    for row in rows:
        row.sort(key=lambda c: c["x"])
    rows.sort(key=lambda r: min(c["y"] for c in r))
    return rows


def attach_connector_labels(children):
    """Attach small floating text shapes to nearby connectors as labels.

    Returns the children list with consumed label shapes removed.
    """
    used = set()
    for c in children:
        if c["kind"] != "connector":
            continue
        best = None
        for t in children:
            if t["kind"] != "text" or id(t) in used:
                continue
            if t["w"] > max(c["w"], c["h"]) * 2 + 60:
                continue
            if c["dir"] == "h":
                cx = (t["x"] + t["x2"]) / 2
                if not (c["x"] - 20 <= cx <= c["x2"] + 20):
                    continue
                dist = min(abs(t["y2"] - c["y"]), abs(c["y2"] - t["y"]))
            else:
                cy = (t["y"] + t["y2"]) / 2
                if not (c["y"] - 20 <= cy <= c["y2"] + 20):
                    continue
                dist = min(abs(t["x2"] - c["x"]), abs(c["x2"] - t["x"]))
            if dist < LABEL_SNAP_PX and (best is None or dist < best[0]):
                best = (dist, t)
        if best:
            used.add(id(best[1]))
            c["label"] = best[1]
    return [s for s in children if id(s) not in used]


# ---------------------------------------------------------------- CSS / HTML

class CssRegistry:
    def __init__(self):
        self.colors, self.rules, self.class_count = {}, {}, {}

    def var(self, hexcolor):
        if not hexcolor:
            return None
        if hexcolor not in self.colors:
            self.colors[hexcolor] = f"--c{len(self.colors) + 1}"
        return f"var({self.colors[hexcolor]})"

    def cls(self, prefix, decls):
        key = (prefix, tuple(sorted(decls.items())))
        if key in self.rules:
            return self.rules[key]
        n = self.class_count.get(prefix, 0) + 1
        self.class_count[prefix] = n
        name = f"{prefix}-{n}" if n > 1 else prefix
        self.rules[key] = name
        return name

    def emit(self):
        lines = []
        for key, name in self.rules.items():
            decls = "; ".join(f"{k}: {v}" for k, v in key[1])
            lines.append(f".{name} {{ {decls}; }}")
        return "\n".join(lines)


def cls_attr(name):
    return f' class="{name}"' if name else ""


def font_stack(name, theme):
    name = name or theme.get("minor_font") or "Calibri"
    return FONT_FALLBACKS.get(name, f"'{name}', sans-serif")


def render_runs(paras, css, theme):
    html = []
    for p in paras:
        if not p["runs"]:
            continue
        pdecl = {}
        if p["align"] and p["align"] != "left":
            pdecl["text-align"] = p["align"]
        if p.get("line_spacing"):
            pdecl["line-height"] = str(p["line_spacing"])
        if p.get("space_after"):
            pdecl["margin-bottom"] = f"{round(p['space_after'] * 1.333, 1)}px"
        pcls = css.cls("para", pdecl) if pdecl else None
        spans = []
        for r in p["runs"]:
            decl = {}
            if r["font"]:
                decl["font-family"] = font_stack(r["font"], theme)
            if r["size"]:
                decl["font-size"] = f"{round(r['size'] * 1.333 / 16, 3)}rem"
            if r["bold"]:
                decl["font-weight"] = "700"
            if r["italic"]:
                decl["font-style"] = "italic"
            if r["color"]:
                decl["color"] = css.var(r["color"])
            scls = css.cls("run", decl) if decl else None
            text = escape(r["text"])
            spans.append(f'<span class="{scls}">{text}</span>' if scls else text)
        body = "".join(spans)
        html.append(f'<p class="{pcls}">{body}</p>' if pcls else f"<p>{body}</p>")
    return "\n".join(html)


def render_table(tbl, css, theme):
    rows_html = []
    for ri, row in enumerate(tbl):
        cells = []
        for cell in row:
            decl = {}
            if cell["fill"]:
                decl["background"] = css.var(cell["fill"])
            ccls = css.cls("cell", decl) if decl else None
            inner = render_runs(cell["paras"], css, theme)
            tag = "th" if ri == 0 else "td"
            cells.append(f"<{tag}{cls_attr(ccls)}>{inner}</{tag}>")
        rows_html.append("<tr>" + "".join(cells) + "</tr>")
    return '<table class="tbl">' + "".join(rows_html) + "</table>"


def shape_box_class(s, css, prefix="panel"):
    decl = {}
    if s["fill"]:
        decl["background"] = css.var(s["fill"])
    if s["border"]:
        decl["border"] = f"{s['border']['width']}px solid {css.var(s['border']['color'])}"
    if s["radius"]:
        decl["border-radius"] = f"{s['radius']}px"
    return css.cls(prefix, decl) if decl else None


def render_divider(s):
    """Plain line with no arrowheads: horizontal rule (vertical lines are
    decorative and don't survive reflow)."""
    if s["dir"] == "v":
        return ""
    return (f'<hr class="divider" style="border-top:{s.get("stroke_w", 1)}px '
            f'solid {s.get("stroke", "#888888")}">')


def render_connector(s, css, theme):
    """Render a line/arrow shape as an inline SVG flow arrow with its label."""
    label = ""
    if s.get("label"):
        label = f'<div class="flow-label">{render_runs(s["label"]["paras"], css, theme)}</div>'
    head = '<path d="M16 2 L4 8 L16 14 Z" stroke="none"/>' if s.get("arrow_start") else ""
    tail = '<path d="M44 2 L56 8 L44 14 Z" stroke="none"/>' if s.get("arrow_end") else ""
    x1 = "12" if s.get("arrow_start") else "2"
    x2 = "48" if s.get("arrow_end") else "58"
    svg = (f'<svg class="arrow-svg" viewBox="0 0 60 16" width="60" height="16" '
           f'fill="currentColor" aria-hidden="true">'
           f'<line x1="{x1}" y1="8" x2="{x2}" y2="8" stroke="currentColor" stroke-width="2.5"/>'
           f"{head}{tail}</svg>")
    vcls = " flow-v" if s["dir"] == "v" else ""
    return (f'<div class="flow-arrow{vcls}" style="color:{s.get("stroke", "#444444")}">'
            f"{label}{svg}</div>")


def render_flow_row(row, cells):
    """A row containing flow arrows: boxes flex, arrows sit between them."""
    parts = []
    for c, h in zip(row, cells):
        parts.append(h if c["kind"] == "connector" else f'<div class="flow-box">{h}</div>')
    return '<div class="flow-row">' + "".join(parts) + "</div>"


def render_node(s, css, theme, depth=0):
    """Recursively render a shape and its children as flex rows/columns."""
    if s["kind"] == "connector":
        return render_connector(s, css, theme)
    if s["kind"] == "divider":
        return render_divider(s)
    if s["kind"] == "img":
        return (f'<img src="{s["img"]}" alt="{escape(s["name"])}" '
                f'style="width:{s["w"]}px;height:{s["h"]}px" class="icon">')
    if s["kind"] == "table":
        return render_table(s["table"], css, theme)

    inner = ""
    if s["paras"]:
        inner = render_runs(s["paras"], css, theme)

    if s["children"]:
        children = attach_connector_labels(s["children"])
        rows = group_rows(children)
        # padding from child offsets
        pad_l = max(0, round(min(c["x"] for c in children) - s["x"]))
        pad_t = max(0, round(min(c["y"] for c in children) - s["y"]))
        pad_r = max(0, round(s["x2"] - max(c["x2"] for c in children)))
        pad_b = max(0, round(s["y2"] - max(c["y2"] for c in children)))
        row_html = []
        prev_bottom = None
        for row in rows:
            top = min(c["y"] for c in row)
            gap_top = round(top - prev_bottom) if prev_bottom is not None else 0
            prev_bottom = max(c["y2"] for c in row)
            cells = [render_node(c, css, theme, depth + 1) for c in row]
            decl = {}
            if gap_top > 2:
                decl["margin-top"] = f"{gap_top}px"
            if any(c["kind"] == "connector" for c in row):
                mt = css.cls("rowgap", decl) if decl else None
                row_html.append(f"<div{cls_attr(mt)}>{render_flow_row(row, cells)}</div>")
                continue
            if len(row) > 1:
                decl["display"] = "flex"
                decl["align-items"] = "center"
                decl["flex-wrap"] = "wrap"
                decl["gap"] = f"{max(8, round(row[1]['x'] - row[0]['x2']))}px"
                rcls = css.cls("row", decl)
                # push the last item right if separated by a wide gap,
                # or if its text is right-aligned (e.g. a header tag)
                last = row[-1]
                right_aligned = bool(last.get("paras")) and any(
                    p["align"] == "right" for p in last["paras"])
                if len(row) >= 2 and (last["x"] - row[-2]["x2"] > 40 or right_aligned):
                    cells[-1] = f'<div class="push-right">{cells[-1]}</div>'
                row_html.append(f'<div class="{rcls}">' + "".join(cells) + "</div>")
            else:
                rcls = css.cls("rowgap", decl) if decl else None
                row_html.append(f"<div{cls_attr(rcls)}>{cells[0]}</div>")
        pad_decl = {"padding": f"{pad_t}px {pad_r}px {pad_b}px {pad_l}px"}
        box = shape_box_class(s, css) or ""
        pcls = css.cls("pad", pad_decl)
        return f'<div class="{(box + " " if box else "")}{pcls}">{inner}{"".join(row_html)}</div>'

    # leaf
    decl = {}
    box = shape_box_class(s, css, "leaf")
    if s["anchor"] and s["h"] > 30 and s["kind"] == "text":
        decl["display"] = "flex"
        decl["flex-direction"] = "column"
        decl["justify-content"] = s["anchor"]
    if s["kind"] == "box" and not s["children"] and not s["paras"]:
        decl["height"] = f"{s['h']}px"   # decorative bar / divider
    extra = css.cls("leafx", decl) if decl else None
    classes = " ".join(c for c in (box, extra) if c)
    return f"<div{cls_attr(classes)}>{inner}</div>"


def convert(pptx_path, out_dir, slide_index=0):
    os.makedirs(out_dir, exist_ok=True)
    assets_dir = os.path.join(out_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    theme = read_theme(prs)
    slide = prs.slides[slide_index]
    layout, master = slide.slide_layout, slide.slide_layout.slide_master
    slide_w = px(prs.slide_width)
    bg = slide_background(slide, theme)
    text_default = scheme_hex("tx1", theme) or "#000000"

    assets = []
    shapes = flatten(slide.shapes, theme, assets_dir, assets, layout, master)
    roots = attach_connector_labels(build_tree(shapes))
    css = CssRegistry()

    rows = group_rows(roots)
    body_rows = []
    prev_bottom = None
    for row in rows:
        top = min(c["y"] for c in row)
        gap = round(top - prev_bottom) if prev_bottom is not None else 0
        prev_bottom = max(c["y2"] for c in row)
        cells = [render_node(c, css, theme) for c in row]
        if any(c["kind"] == "connector" for c in row):
            style = f' style="margin-top:{gap}px"' if gap > 2 else ""
            body_rows.append(f"<div{style}>{render_flow_row(row, cells)}</div>")
        elif len(row) > 1:
            gap_x = max(16, round(row[1]["x"] - row[0]["x2"]))
            widths = [max(c["w"], 40) for c in row]
            if max(widths) / min(widths) > 1.5 or len(row) > 3:
                gtc = " ".join(f"{round(w)}fr" for w in widths)
                grid_style = f"grid-template-columns:{gtc}; --gap:{gap_x}px"
            else:
                grid_style = f"--gap:{gap_x}px"
            # small inline items center nicely; tall cards align to the top
            align = "center" if max(c["h"] for c in row) < 160 else "start"
            grid_style += f"; align-items:{align}"
            body_rows.append(
                f'<div class="grid-row" style="margin-top:{max(gap, 0)}px; {grid_style}">'
                + "".join(f'<div class="grid-cell">{c}</div>' for c in cells) + "</div>")
        else:
            style = f' style="margin-top:{gap}px"' if gap > 2 else ""
            body_rows.append(f"<div{style}>{cells[0]}</div>")

    side_margins = [s["x"] for s in roots if s["w"] < slide_w * 0.99]
    margin_l = round(min(side_margins)) if side_margins else 48

    color_vars = "\n".join(f"  {v}: {k};" for k, v in css.colors.items())
    stylesheet = f""":root {{
{color_vars}
  --bg: {bg};
  --text: {text_default};
  --font-body: {font_stack(theme.get('minor_font'), theme)};
  --page-max: {round(slide_w)}px;
}}
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
html {{ font-size: 16px; }}
body {{ font-family: var(--font-body); background: var(--bg); color: var(--text); }}
.dashboard {{ max-width: var(--page-max); margin: 0 auto; padding: 0 {margin_l}px 48px; }}
.full-bleed {{ margin: 0 -{margin_l}px; }}
p {{ line-height: 1.4; }}
.icon {{ flex: 0 0 auto; object-fit: contain; }}
.push-right {{ margin-left: auto; }}
.divider {{ border: none; width: 100%; }}
.grid-row {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(320px, 1fr)); gap: var(--gap, 24px); }}
.flow-row {{ display: flex; align-items: stretch; gap: 16px; }}
.flow-row .flow-box {{ flex: 1 1 280px; min-width: 0; }}
.flow-arrow {{ flex: 0 0 auto; display: flex; flex-direction: column; align-items: center;
  justify-content: center; gap: 4px; align-self: center; }}
.flow-arrow .flow-label {{ white-space: nowrap; }}
.flow-arrow.flow-v .arrow-svg {{ transform: rotate(90deg); }}
.tbl {{ border-collapse: collapse; width: 100%; }}
.tbl th, .tbl td {{ padding: 8px 12px; text-align: left; }}
{css.emit()}
@media (max-width: 640px) {{
  html {{ font-size: 14px; }}
  .dashboard {{ padding-left: 20px; padding-right: 20px; }}
  .full-bleed {{ margin: 0 -20px; }}
  .grid-row {{ grid-template-columns: 1fr !important; }}
  .flow-row {{ flex-direction: column; }}
  .flow-row .flow-box {{ flex: 1 1 auto; width: 100%; }}
  .flow-arrow .arrow-svg {{ transform: rotate(90deg); margin: 16px 0; }}
}}
"""
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{escape(os.path.splitext(os.path.basename(pptx_path))[0])}</title>
<link rel="stylesheet" href="styles.css">
</head>
<body>
<main class="dashboard">
{chr(10).join(body_rows)}
</main>
</body>
</html>
"""
    with open(os.path.join(out_dir, "index.html"), "w", encoding="utf-8") as fh:
        fh.write(html)
    with open(os.path.join(out_dir, "styles.css"), "w", encoding="utf-8") as fh:
        fh.write(stylesheet)
    print(f"Wrote {out_dir}/index.html, styles.css, {len(assets)} asset(s)")


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx")
    ap.add_argument("-o", "--out", default="html_out")
    ap.add_argument("--slide", type=int, default=1, help="1-based slide number")
    args = ap.parse_args()
    convert(args.pptx, args.out, args.slide - 1)


if __name__ == "__main__":
    main()
